#!/usr/bin/env python3
import os
import shlex
import subprocess
from pptx import Presentation
from pptx.util import Inches


ROOT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DIAG_DIR = os.path.join(ROOT_DIR, "assets", "diagrams")
OUT_DIR = os.path.join(ROOT_DIR, "dist")
os.makedirs(OUT_DIR, exist_ok=True)


slides = [
    # title, bullets, image (optional), notes
    (
        "Spec-Driven Development – KI-fähige Softwareentwicklung",
        ["Prozesse, Rollen und Artefakte für echte Produktivität"],
        None,
        "Fokus auf Strukturänderungen für verlässliche KI-Nutzung.",
    ),
    (
        "Executive Summary",
        [
            "Agilität: Mensch=Kommunikation, KI=Explizitheit",
            "Specs sind Quelle der Wahrheit; Review nach vorne",
            "KI skaliert Tempo und Fehler – Entscheidungsqualität wird Engpass",
        ],
        None,
        "Thesen, die den Rest des Vortrags rahmen.",
    ),
    (
        "Warum Wandel: Von Code- zu Prozessfokus",
        [
            "Traditionell: Feature→Ticket→PR",
            "Prozessbasiert: Idee→Spec→Gen→Gates→Review",
            "70–80% mechanisch automatisierbar",
        ],
        None,
        "Vom Tippen zu System- und Prozessdesign.",
    ),
    (
        "Agilität heute (Kommunikation als Klebstoff)",
        [],
        os.path.join(DIAG_DIR, "agility_today.png"),
        "Menschen gleichen Unschärfen durch Gespräche aus.",
    ),
    (
        "Zwei Agilitäten: Human↔Human / Human↔Machine",
        [],
        os.path.join(DIAG_DIR, "dual_agility.png"),
        "Gespräche enden in Spezifikationen; Stakeholder entscheiden.",
    ),
    (
        "Implizit → Explizit: Der große Shift",
        [
            "Architektur formulieren, Regeln schreiben",
            "Wissen versionieren, Artefakte binden",
            "Git=Was, Confluence=Warum, Jira=Steuerung",
        ],
        None,
        "Specs als formalisierte Kommunikation Mensch↔Maschine.",
    ),
    (
        "SDD: Prinzipien",
        [
            "Specs=Source of Truth; Code implementiert Specs",
            "Verhaltensänderung → Spec-Änderung",
            "Lokale Gates (Format/Lint/Tests)",
        ],
        None,
        "Ziel: günstige Evolution, nicht Bürokratie.",
    ),
    (
        "Wasserfall vs. SDD (Abgrenzung)",
        [],
        os.path.join(DIAG_DIR, "waterfall_vs_sdd.png"),
        "Spec zuerst und ständig; Feedback früh, Ableitung billig.",
    ),
    (
        "Tests = ausführbare Specs",
        [
            "Use-Case/Command-Tests (Black-Box)",
            "ACs beweisen Verhalten",
            "Architektur & Verträge schützen",
        ],
        None,
        "DoD: ACs grün.",
    ),
    (
        "KI ist Umsetzer, nicht Teilnehmer",
        ["Kein Teammitglied/Stakeholder", "Umsetzer, Varianten-Generator, Verstärker"],
        None,
        "Diskussionen menschlich; Umsetzung maschinell.",
    ),
    (
        "SDD-Regelkreis",
        [],
        os.path.join(DIAG_DIR, "sdd_loop.png"),
        "Conversation → Decision → Spec → Generate → Learn.",
    ),
    (
        "Multi-Agent-Workflow (lokal)",
        ["Local-only; Rollen/Output-Verträge", "Deterministische Gates; Artefakt-Handoffs"],
        None,
        "Reproduzierbar & auditierbar, ohne SaaS.",
    ),
    (
        "Repo-Struktur, Agenten & Handoffs",
        [
            "specs/, .github/agents, scripts/gates",
            "Architect/Planner/Implementer/Tester/Reviewer",
            "AUTOMATION MODE = diff-only",
        ],
        None,
        "Klare Zuständigkeiten und Übergaben.",
    ),
    (
        "Gates & Orchestrator",
        ["Format/Lint/Tests lokal", "State-Machine-Flow", "Kleine, reviewbare Patches"],
        None,
        "Fehler minimal & lokal fixen; Loops bis grün.",
    ),
    (
        "Review-Shift & Rituale",
        ["Spec→Gen→Diff→Code", "Dailies/Retros menschlich", "Refinement → Spec/Tasks"],
        None,
        "Code-Reviews ohne Specs sind Meinungen.",
    ),
    (
        "CI/CD als Design-Pipeline",
        [],
        os.path.join(DIAG_DIR, "ci_design_pipeline.png"),
        "Varianten/Previews; ‘Show, don’t tell’.",
    ),
    (
        "MVP als Verhaltensschnitt",
        ["Happy Path zuerst", "Specs vollständig, Implementierung inkrementell"],
        None,
        "Frühe Validierung von Schnittstellen/Regeln.",
    ),
    (
        "Brownfield: SpecKit > BMAD",
        ["Mini-Specs/Constraints/Tasks", "BMAD zu schwer für Brownfield"],
        None,
        "‘Fix X, brich Y nicht’ mit klaren Tests.",
    ),
    (
        "Anti-Patterns & Risiken",
        ["Prompt-to-Feature ohne Spec", "Confluence/Jira als Wahrheit", "Spec-Freeze/Kultur"],
        None,
        "Overhead kurzfr.; Rework sinkt mittel-/langfr.",
    ),
    (
        "30-Tage-Plan + Metriken",
        ["Woche 1–4 Schritte", "Zeit bis Preview; Spec-Änderungen; Rework; Gate-Stabilität"],
        None,
        "Start klein mit End-to-End Slice.",
    ),
]


def _run_mermaid(mmd_path: str, png_path: str) -> None:
    # Prefer installed mmdc; fallback to npx.
    try:
        subprocess.run(
            ["mmdc", "-i", mmd_path, "-o", png_path, "-b", "transparent"],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        return
    except Exception:
        pass

    cmd = f"npx --yes @mermaid-js/mermaid-cli -i {shlex.quote(mmd_path)} -o {shlex.quote(png_path)} -b transparent"
    subprocess.run(cmd, shell=True, check=True)


def add_title(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bullets(
    prs: Presentation, title: str, bullets: list[str], notes: str | None, image: str | None = None
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    if image:
        slide.shapes.add_picture(image, Inches(0.8), Inches(1.8), width=Inches(8.4))
    else:
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(bullets):
            paragraph = body.add_paragraph() if i > 0 else body.paragraphs[0]
            paragraph.text = bullet
            paragraph.level = 0
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def main() -> None:
    diagrams = [
        "agility_today",
        "dual_agility",
        "waterfall_vs_sdd",
        "sdd_loop",
        "ci_design_pipeline",
    ]

    for name in diagrams:
        mmd = os.path.join(DIAG_DIR, f"{name}.mmd")
        png = os.path.join(DIAG_DIR, f"{name}.png")
        if os.path.exists(mmd) and not os.path.exists(png):
            _run_mermaid(mmd, png)

    prs = Presentation()
    add_title(prs, slides[0][0], slides[0][1][0])
    for title, bullets, image, notes in slides[1:]:
        add_bullets(prs, title, bullets, notes, image=image)

    out = os.path.join(OUT_DIR, "Spec-Driven_Engineering_AI-fähig.pptx")
    prs.save(out)
    print(f"✅ Saved: {out}")


if __name__ == "__main__":
    main()
